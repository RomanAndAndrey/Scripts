"""
Генератор презентаций PowerPoint из Markdown файла.
Читает PRESENTATION_PLAN.md и создает Defense.pptx.
"""

import os
import re
from typing import List, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.slide import Slide
from pptx.util import Inches, Pt


def main() -> None:
    """
    Главная функция генерации презентации.
    Читает PRESENTATION_PLAN.md и создает Defense.pptx.
    """
    script_dir = os.path.dirname(os.path.abspath(__file__))
    plan_path = os.path.join(script_dir, "PRESENTATION_PLAN.md")
    output_path = os.path.join(script_dir, "Defense.pptx")
    project_root = os.path.dirname(script_dir)

    print(f"Reading plan from: {plan_path}")
    if not os.path.exists(plan_path):
        print("Plan file not found!")
        return

    # Инициализация презентации
    prs = Presentation()

    with open(plan_path, "r", encoding="utf-8") as f:
        lines = f.readlines()

    # Состояние парсера
    state: str = "NORMAL"  # NORMAL, TITLE, CONTENT, ARCH, CODE_SECTION, IN_CODE_BLOCK
    current_slide: Optional[Slide] = None
    code_buffer: List[str] = []

    for line in lines:
        raw_line = line
        line = line.strip()

        # Обработка заголовков уровня 2 (##)
        if line.startswith("## "):
            current_slide, state = process_header(prs, line, state, code_buffer, current_slide)
            code_buffer.clear()
            continue

        # Обработка подзаголовков уровня 3 (###) для слайдов с кодом
        if line.startswith("### ") and state in ["CODE_SECTION", "IN_CODE_BLOCK", "NORMAL"]:
            sub_header = line.replace("###", "").strip()
            current_slide = prs.slides.add_slide(prs.slide_layouts[5])
            current_slide.shapes.title.text = sub_header
            state = "CODE_WAITING"
            continue

        # Обработка блоков кода
        if line.startswith("```"):
            if state == "IN_CODE_BLOCK":
                # Конец блока кода
                add_code_to_slide(current_slide, code_buffer)
                code_buffer.clear()
                state = "CODE_SECTION"
            else:
                # Начало блока кода
                state = "IN_CODE_BLOCK"
                code_buffer.clear()
            continue

        if state == "IN_CODE_BLOCK":
            code_buffer.append(raw_line)
            continue

        # Обработка содержимого слайдов
        process_content(current_slide, line, state, project_root, script_dir)

    prs.save(output_path)
    print(f"Successfully generated: {output_path}")


def process_header(
    prs: Presentation, line: str, state: str, code_buffer: List[str], current_slide: Optional[Slide]
) -> tuple[Optional[Slide], str]:
    """
    Обрабатывает заголовки уровня 2 и создает слайды.

    Args:
        prs: Объект презентации
        line: Строка с заголовком
        state: Текущее состояние парсера
        code_buffer: Буфер с кодом
        current_slide: Текущий слайд

    Returns:
        Кортеж (новый слайд, новое состояние)
    """
    header_text = line.replace("##", "").strip()
    header_clean = re.sub(r"^\d+\.\s*", "", header_text)

    print(f"Processing Section: {header_clean}")

    if "Титульный" in header_text or "Title" in header_text:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        if slide.shapes.title:
            slide.shapes.title.text = ""
        if slide.placeholders[1]:
            slide.placeholders[1].text = ""
        return slide, "TITLE"

    elif "Архитектура" in header_text:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = header_clean
        return slide, "ARCH"

    elif "Ключевые фрагменты" in header_text or "Код" in header_text:
        return current_slide, "CODE_SECTION"

    else:
        # Стандартный слайд с буллетами
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = header_clean
        if slide.placeholders[1].has_text_frame:
            slide.placeholders[1].text_frame.clear()
        return slide, "CONTENT"


def add_code_to_slide(slide: Optional[Slide], code_buffer: List[str]) -> None:
    """
    Добавляет блок кода на слайд.

    Args:
        slide: Слайд для добавления кода
        code_buffer: Список строк с кодом
    """
    if not slide:
        return

    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(5)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    text_content = "".join(code_buffer)
    p = tf.add_paragraph()
    p.text = text_content
    p.font.name = "Consolas"
    p.font.size = Pt(11)


def process_content(
    slide: Optional[Slide], line: str, state: str, project_root: str, script_dir: str
) -> None:
    """
    Обрабатывает содержимое слайдов в зависимости от состояния.

    Args:
        slide: Текущий слайд
        line: Строка для обработки
        state: Текущее состояние парсера
        project_root: Корневая директория проекта
        script_dir: Директория скрипта
    """
    if not slide:
        return

    if state == "TITLE" and line.startswith("-"):
        # Обработка титульного слайда
        clean_line = line.lstrip("- ").replace("**", "")
        if ":" in clean_line:
            parts = clean_line.split(":", 1)
            key = parts[0].strip()
            val = parts[1].strip()

            if "Название" in key or "Project" in key:
                slide.shapes.title.text = val
            else:
                # Подзаголовок
                ph = slide.placeholders[1]
                p = ph.text_frame.add_paragraph()
                p.text = val
                p.font.name = "Calibri"
                p.font.size = Pt(24)

    elif state == "CONTENT" and line.startswith("-"):
        # Добавление буллета
        text = line.lstrip("- ").replace("**", "").strip()
        if slide.placeholders[1].has_text_frame:
            tf = slide.placeholders[1].text_frame
            p = tf.add_paragraph()
            p.text = text
            p.level = 0
            p.font.name = "Calibri"
            p.font.color.rgb = RGBColor(0, 0, 0)

    elif state == "ARCH":
        # Обработка архитектурной диаграммы
        if "db_schema.png" in line:
            add_architecture_image(slide, project_root, script_dir)


def add_architecture_image(slide: Slide, project_root: str, script_dir: str) -> None:
    """
    Добавляет изображение архитектуры на слайд.

    Args:
        slide: Слайд для добавления изображения
        project_root: Корневая директория проекта
        script_dir: Директория скрипта
    """
    img_name = "db_schema.png"
    paths_to_check = [os.path.join(project_root, img_name), os.path.join(script_dir, img_name)]

    found = False
    for p in paths_to_check:
        if os.path.exists(p):
            print(f"Embedding image: {p}")
            slide.shapes.add_picture(p, Inches(1), Inches(1.5), height=Inches(5))
            found = True
            break

    if not found:
        print("Image not found, adding placeholder.")
        tx = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
        tx.text_frame.text = "[Изображение db_schema.png не найдено]"
        tx.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


if __name__ == "__main__":
    main()
